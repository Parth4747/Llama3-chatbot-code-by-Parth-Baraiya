import PyPDF2
import chainlit as cl
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ChatMessageHistory, ConversationBufferMemory
from langchain_groq import ChatGroq
from dotenv import load_dotenv
import os
from datetime import datetime

# Loading environment variables from .env file
load_dotenv() 

# Function to initialize conversation chain with GROQ language model
groq_api_key = os.environ['GROQ_API_KEY']

# Initializing GROQ chat with provided API key, model name, and settings
llm_groq = ChatGroq(
            groq_api_key=groq_api_key, model_name="llama3-70b-8192",
            temperature=0.2
)

# Global variables for files and chain
uploaded_files = None
chain = None

@cl.on_chat_start
async def on_chat_start():
    # Prompt the user to select an action
    options = ["Scanning", "Chat"]
    await cl.Message(content="Welcome! Please type 'Scanning' to upload and scan your files to start asking questions.").send()

@cl.on_message
async def on_message(message: cl.Message):
    global uploaded_files, chain

    if message.content.strip().lower() == "scanning":
        files = None  # Initialize variable to store uploaded files

        # Wait for the user to upload files
        while files is None:
            files = await cl.AskFileMessage(
                content=f"{datetime.now().strftime('%I:%M:%S %p')}\nPlease upload one or more files to begin!",
                accept=["application/pdf", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/vnd.openxmlformats-officedocument.presentationml.presentation", "text/csv", "text/plain"],
                max_size_mb=100,
                max_files=10,
                timeout=180,  # Set a timeout for user response
            ).send()

        uploaded_files = files  # Store the uploaded files globally

        # Process each uploaded file
        texts = []
        metadatas = []
        for file in files:
            print(file)  # Print the file object for debugging

            # Read the file content based on file type
            if file.type == "application/pdf":
                pdf = PyPDF2.PdfReader(file.path)
                file_text = ""
                for page in pdf.pages:
                    file_text += page.extract_text()
            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                # Process .docx files
                from docx import Document
                doc = Document(file.path)
                file_text = "\n".join([p.text for p in doc.paragraphs])
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                # Process .ppt files
                import pptx
                ppt = pptx.Presentation(file.path)
                file_text = ""
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            file_text += shape.text
                        elif shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    file_text += run.text
                        file_text += "\n"
            elif file.type == "text/csv" or file.type == "text/plain":
                # Process .csv or .txt files
                with open(file.path, 'r') as f:
                    file_text = f.read()

            # Split the text into chunks
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1200, chunk_overlap=50)
            file_texts = text_splitter.split_text(file_text)
            texts.extend(file_texts)

            # Create a metadata for each chunk
            file_metadatas = [{"source": f"{i}-{file.name}"} for i in range(len(file_texts))]
            metadatas.extend(file_metadatas)

        # Create a Chroma vector store
        embeddings = OllamaEmbeddings(model="nomic-embed-text")
        docsearch = await cl.make_async(Chroma.from_texts)(
            texts, embeddings, metadatas=metadatas
        )
        
        # Initialize message history for conversation
        message_history = ChatMessageHistory()
        
        # Memory for conversational context
        memory = ConversationBufferMemory(
            memory_key="chat_history",
            output_key="answer",
            chat_memory=message_history,
            return_messages=True,
        )

        # Create a chain that uses the Chroma vector store
        chain = ConversationalRetrievalChain.from_llm(
            llm=llm_groq,
            chain_type="stuff",
            retriever=docsearch.as_retriever(),
            memory=memory,
            return_source_documents=True,
        )
        
        # Inform the user that processing has ended. You can now chat.
        timestamp = datetime.now().strftime('%I:%M:%S %p')
        await cl.Message(content=f"{timestamp}\nProcessing {len(files)} files done. You can now ask questions!").send()

    elif message.content.strip().lower() == "chat":
        if chain is None:
            await cl.Message(content="Please upload and scan your files first by typing 'Scanning'.").send()
            return

        await cl.Message(content="You can now ask your questions.").send()

    else:
        if chain is None:
            await cl.Message(content="Please upload and scan your files first by typing 'Scanning'.").send()
            return

        # Callbacks happen asynchronously/parallel 
        cb = cl.AsyncLangchainCallbackHandler()
        
        # Call the chain with user's message content
        res = await chain.ainvoke(message.content, callbacks=[cb])
        answer = res["answer"]
        source_documents = res["source_documents"] 

        text_elements = []  # Initialize list to store text elements
        
        # Process source documents if available
        if source_documents:
            for source_idx, source_doc in enumerate(source_documents):
                source_name = f"source_{source_idx}"
                # Create the text element referenced in the message
                text_elements.append(
                    cl.Text(content=source_doc.page_content, name=source_name)
                )
            source_names = [text_el.name for text_el in text_elements]
            
            # Add source references to the answer with timestamp
            if source_names:
                answer += f"\nSources: {', '.join(source_names)}\nTimestamp: {datetime.now().strftime('%I:%M:%S %p')}"
            else:
                answer += f"\nNo sources found\nTimestamp: {datetime.now().strftime('%I:%M:%S %p')}"

        # Return results with current timestamp
        timestamp = datetime.now().strftime('%I:%M:%S %p')
        await cl.Message(content=f"{timestamp}\n{answer}", elements=text_elements).send()
